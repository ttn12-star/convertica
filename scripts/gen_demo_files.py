#!/usr/bin/env python3
"""Generate the demo files used as screenshot props (scripts/demo_files/).

Filenames are deliberate easter eggs (song/anime/film TITLES only — titles
are not copyrightable; no lyrics, no third-party artwork). PDF/image contents
are original: tools like organize/sign render page thumbnails client-side,
so those pages are designed to look good in screenshots.

    python scripts/gen_demo_files.py
"""
from __future__ import annotations

import io
import zipfile
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from reportlab.lib.colors import HexColor
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas

OUT = Path(__file__).resolve().parent / "demo_files"
OUT.mkdir(exist_ok=True)

INDIGO = HexColor("#4f46e5")
VIOLET = HexColor("#7c3aed")
INK = HexColor("#1e1b4b")
PAPER = HexColor("#f8fafc")


def _page_frame(c, title: str, accent):
    w, h = A4
    c.setFillColor(PAPER)
    c.rect(0, 0, w, h, stroke=0, fill=1)
    c.setFillColor(accent)
    c.rect(0, h - 16, w, 16, stroke=0, fill=1)
    c.setFillColor(INK)
    c.setFont("Helvetica-Bold", 22)
    c.drawString(56, h - 90, title)
    c.setStrokeColor(accent)
    c.setLineWidth(2)
    c.line(56, h - 104, w - 56, h - 104)


def _body_lines(c, lines, y_start=None):
    w, h = A4
    y = y_start or h - 150
    c.setFont("Helvetica", 12)
    c.setFillColor(HexColor("#334155"))
    for line in lines:
        c.drawString(56, y, line)
        y -= 22


def pdf_titan_plan():
    path = OUT / "plan_to_defeat_titans_v2_FINAL.pdf"
    c = canvas.Canvas(str(path), pagesize=A4)
    _page_frame(c, "Plan to Defeat Titans — v2 (FINAL, really)", INDIGO)
    _body_lines(
        c,
        [
            "Phase 1: Reclaim the outer wall (budget: 3 coffees)",
            "Phase 2: Very tall ladders",
            "Phase 3: Ask them politely to leave",
            "",
            "Risk register: see appendix, page 67.",
        ],
    )
    c.setFont("Helvetica-Bold", 10)
    c.setFillColor(VIOLET)
    c.drawString(56, 60, "CONFIDENTIAL-ISH  ·  copy 42 of 67")
    c.showPage()
    _page_frame(c, "Appendix A: Wall Maintenance Schedule", VIOLET)
    _body_lines(c, ["Monday: patrol", "Tuesday: patrol", "Wednesday: long patrol"])
    c.showPage()
    c.save()
    return path


def pdf_chapters(n_pages: int = 8):
    """Big numbered pages — looks great in organize/split thumbnail grids."""
    path = OUT / "chapter_67_storyboard.pdf"
    c = canvas.Canvas(str(path), pagesize=A4)
    palette = [
        "#4f46e5",
        "#7c3aed",
        "#2563eb",
        "#9333ea",
        "#4338ca",
        "#6d28d9",
        "#3b82f6",
        "#8b5cf6",
    ]
    w, h = A4
    for i in range(n_pages):
        accent = HexColor(palette[i % len(palette)])
        c.setFillColor(PAPER)
        c.rect(0, 0, w, h, stroke=0, fill=1)
        c.setFillColor(accent)
        c.circle(w / 2, h / 2 + 40, 110, stroke=0, fill=1)
        c.setFillColor(HexColor("#ffffff"))
        c.setFont("Helvetica-Bold", 96)
        c.drawCentredString(w / 2, h / 2 + 8, str(i + 1))
        c.setFillColor(INK)
        c.setFont("Helvetica-Bold", 16)
        c.drawCentredString(w / 2, 120, f"Storyboard — scene {i + 1} of 67")
        c.showPage()
    c.save()
    return path


def pdf_chords():
    path = OUT / "hallelujah_chords.pdf"
    c = canvas.Canvas(str(path), pagesize=A4)
    _page_frame(c, "Hallelujah — chords (practice sheet)", VIOLET)
    _body_lines(
        c,
        [
            "Intro:  C  ·  Am  ·  C  ·  Am",
            "Verse:  C  ·  Am  ·  F  ·  G  ·  E7  ·  Am",
            "",
            "Note to self: practice the 4/4 part slowly.",
            "(Lyrics not included — you know them anyway.)",
        ],
    )
    c.showPage()
    c.save()
    return path


def docx_rickroll():
    from docx import Document

    path = OUT / "never_gonna_give_you_up.docx"
    d = Document()
    d.add_heading("Never Gonna Give You Up", level=1)
    d.add_paragraph("A short manifesto about commitment to your files.")
    d.add_paragraph(
        "We are no strangers to conversions. You know the rules, and so do we: "
        "no file left behind, no format left unsupported."
    )
    d.save(str(path))
    return path


def xlsx_totoro():
    from openpyxl import Workbook

    path = OUT / "totoro_watchlist.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "Watchlist"
    ws.append(["Title", "Mood", "Rewatch?"])
    for row in [
        ("My Neighbor Totoro", "cozy", "always"),
        ("Spirited Away", "wondrous", "yes"),
        ("Porco Rosso", "aviation", "yes"),
        ("Whisper of the Heart", "soft", "67 times"),
    ]:
        ws.append(row)
    wb.save(str(path))
    return path


def pptx_daft():
    from pptx import Presentation

    path = OUT / "one_more_time_setlist.pptx"
    p = Presentation()
    slide = p.slides.add_slide(p.slide_layouts[0])
    slide.shapes.title.text = "One More Time — setlist"
    slide.placeholders[1].text = (
        "Track 1: intro · Track 2: the good one · Encore: same but louder"
    )
    p.save(str(path))
    return path


def _doodle_image(size, bg1, bg2, caption):
    """Original doodle art: gradient + rubber duck + caption."""
    w, h = size
    img = Image.new("RGB", size)
    dr = ImageDraw.Draw(img)
    for y in range(h):
        t = y / h
        rgb = tuple(int(a + (b - a) * t) for a, b in zip(bg1, bg2, strict=False))
        dr.line([(0, y), (w, y)], fill=rgb)
    # rubber duck silhouette (original, blobby on purpose)
    cx, cy, r = w * 0.5, h * 0.58, min(w, h) * 0.18
    dr.ellipse(
        [cx - r * 1.4, cy - r * 0.6, cx + r * 1.4, cy + r * 0.9], fill=(255, 214, 10)
    )
    dr.ellipse(
        [cx + r * 0.6, cy - r * 1.5, cx + r * 1.8, cy - r * 0.3], fill=(255, 214, 10)
    )
    dr.ellipse(
        [cx + r * 1.5, cy - r * 1.05, cx + r * 1.95, cy - r * 0.75], fill=(255, 122, 0)
    )
    dr.ellipse(
        [cx + r * 1.15, cy - r * 1.2, cx + r * 1.3, cy - r * 1.05], fill=(30, 27, 75)
    )
    try:
        font = ImageFont.truetype(
            "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", int(h * 0.07)
        )
    except OSError:
        font = ImageFont.load_default()
    dr.text((w * 0.5, h * 0.88), caption, fill=(255, 255, 255), font=font, anchor="mm")
    return img


def images():
    paths = []
    duck = _doodle_image(
        (1200, 900), (79, 70, 229), (124, 58, 237), "much wow · very convert"
    )
    p = OUT / "much_wow_vacation_photo.jpg"
    duck.save(p, quality=88)
    paths.append(p)

    heart = _doodle_image(
        (800, 800), (37, 99, 235), (147, 51, 234), "pixel heart says hi"
    )
    p = OUT / "pixel_heart_collection.png"
    heart.save(p)
    paths.append(p)

    webp = _doodle_image((900, 700), (67, 56, 202), (59, 130, 246), "level 67 unlocked")
    p = OUT / "level_67_unlocked.webp"
    webp.save(p, "WEBP", quality=85)
    paths.append(p)

    ico = _doodle_image((256, 256), (79, 70, 229), (124, 58, 237), "")
    p = OUT / "not_a_moon_favicon.ico"
    ico.save(p, sizes=[(16, 16), (32, 32), (48, 48), (64, 64), (128, 128), (256, 256)])
    paths.append(p)

    # HEIC can't be written by Pillow; the picker only shows the name, the
    # file is never converted during capture — a tiny stub is enough.
    p = OUT / "skywalker_family_album.heic"
    p.write_bytes(b"\x00\x00\x00\x18ftypheic" + b"\x00" * 64)
    paths.append(p)
    return paths


def misc_text():
    paths = []
    p = OUT / "legendary_sword_logo.svg"
    p.write_text(
        '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100">'
        '<defs><linearGradient id="g" x1="0" y1="0" x2="1" y2="1">'
        '<stop offset="0" stop-color="#4f46e5"/><stop offset="1" stop-color="#7c3aed"/>'
        "</linearGradient></defs>"
        '<rect width="100" height="100" rx="20" fill="url(#g)"/>'
        '<path d="M50 15 54 60 50 70 46 60Z" fill="#fff"/>'
        '<rect x="40" y="60" width="20" height="6" rx="3" fill="#fbbf24"/>'
        "</svg>"
    )
    paths.append(p)

    p = OUT / "hello_there_kenobi.html"
    p.write_text(
        "<!doctype html><html><head><title>Hello there</title></head>"
        "<body style='font-family:sans-serif'><h1>Hello there.</h1>"
        "<p>General greeting page. Perfectly balanced markup.</p></body></html>"
    )
    paths.append(p)

    p = OUT / "shrek_screenplay_draft_9.txt"
    p.write_text(
        "SHREK-INSPIRED ORIGINAL DRAFT 9\n\n"
        "Scene 1. A swamp. Someone really values privacy.\n"
        "Scene 2. Layers are discussed. Like onions. And PDFs.\n"
    )
    paths.append(p)

    p = OUT / "README_secret_lore.md"
    p.write_text(
        "# Secret Lore\n\n- there are exactly 67 easter eggs\n- or are there?\n"
    )
    paths.append(p)
    return paths


def epub_and_zip(pdf_for_zip: Path):
    paths = []
    p = OUT / "fellowship_travel_notes.epub"
    with zipfile.ZipFile(p, "w") as z:
        z.writestr("mimetype", "application/epub+zip")
        z.writestr(
            "META-INF/container.xml",
            '<?xml version="1.0"?><container version="1.0" '
            'xmlns="urn:oasis:names:tc:opendocument:xmlns:container">'
            '<rootfiles><rootfile full-path="content.opf" '
            'media-type="application/oebps-package+xml"/></rootfiles></container>',
        )
        z.writestr(
            "content.opf",
            '<?xml version="1.0"?><package xmlns="http://www.idpf.org/2007/opf" version="3.0" '
            'unique-identifier="id"><metadata xmlns:dc="http://purl.org/dc/elements/1.1/">'
            '<dc:title>Fellowship Travel Notes</dc:title><dc:identifier id="id">demo-67</dc:identifier>'
            "<dc:language>en</dc:language></metadata>"
            '<manifest><item id="c" href="ch1.xhtml" media-type="application/xhtml+xml"/></manifest>'
            '<spine><itemref idref="c"/></spine></package>',
        )
        z.writestr(
            "ch1.xhtml",
            '<?xml version="1.0"?><html xmlns="http://www.w3.org/1999/xhtml"><head>'
            "<title>Ch 1</title></head><body><h1>Simply walking. Chapter 1 of 67.</h1></body></html>",
        )
    paths.append(p)

    p = OUT / "winter_is_coming_backup.zip"
    with zipfile.ZipFile(p, "w") as z:
        z.write(pdf_for_zip, pdf_for_zip.name)
        z.writestr("house_words_draft.txt", "Winter is coming. Back up your files.\n")
    paths.append(p)
    return paths


if __name__ == "__main__":
    made = [
        pdf_titan_plan(),
        pdf_chapters(),
        pdf_chords(),
        docx_rickroll(),
        xlsx_totoro(),
        pptx_daft(),
        *images(),
        *misc_text(),
    ]
    made += epub_and_zip(made[2])
    for p in made:
        assert p.stat().st_size > 0, p
        print(f"{p.name}  {p.stat().st_size // 1024}K")
    print(f"OK: {len(made)} demo files in {OUT}")
