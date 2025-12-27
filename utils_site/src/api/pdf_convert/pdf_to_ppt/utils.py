"""
PDF to PowerPoint conversion utilities.

Note: This is a basic implementation that extracts PDF pages as images
and creates a PowerPoint presentation. For production use, consider using
specialized libraries or services for better quality conversion.
"""

import os
import tempfile
from pathlib import Path

from django.core.files.uploadedfile import UploadedFile
from django.utils.text import get_valid_filename
from src.api.file_validation import check_disk_space, sanitize_filename
from src.api.logging_utils import get_logger
from src.exceptions import ConversionError, StorageError

logger = get_logger(__name__)


def convert_pdf_to_ppt(
    uploaded_file: UploadedFile,
    extract_images: bool = True,
    suffix: str = "_convertica",
) -> tuple[str, str]:
    """Convert PDF to PowerPoint presentation.

    Args:
        uploaded_file: PDF file to convert
        extract_images: Whether to extract images from PDF
        suffix: Suffix for output filename

    Returns:
        Tuple of (input_path, output_path)

    Raises:
        ConversionError: If conversion fails
        StorageError: If disk space is insufficient
    """
    context = {
        "function": "convert_pdf_to_ppt",
        "input_filename": os.path.basename(uploaded_file.name),
        "input_size": uploaded_file.size,
        "extract_images": extract_images,
    }

    logger.info("Starting PDF to PowerPoint conversion", extra=context)

    try:
        from pdf2image import convert_from_path
    except ModuleNotFoundError as e:
        raise ConversionError(
            "PDF to PowerPoint conversion requires 'pdf2image' to be installed."
        ) from e

    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ModuleNotFoundError as e:
        raise ConversionError(
            "PDF to PowerPoint conversion requires 'python-pptx' to be installed."
        ) from e

    # Create temp directory
    tmp_dir = tempfile.mkdtemp(prefix="pdf_to_ppt_")
    input_path = None
    output_path = None

    try:
        # Check disk space
        required_space = uploaded_file.size * 5  # Estimate
        check_disk_space(required_space, context)

        # Save uploaded file
        safe_filename = sanitize_filename(get_valid_filename(uploaded_file.name))
        input_path = os.path.join(tmp_dir, safe_filename)

        with open(input_path, "wb") as f:
            for chunk in uploaded_file.chunks():
                f.write(chunk)

        logger.debug(
            "Saved PDF file",
            extra={
                **context,
                "input_path": input_path,
                "file_size": uploaded_file.size,
            },
        )

        # Convert PDF pages to images
        logger.info("Converting PDF pages to images", extra=context)
        images = convert_from_path(input_path, dpi=150)

        if not images:
            raise ConversionError("Failed to extract pages from PDF")

        logger.info(
            f"Extracted {len(images)} pages from PDF",
            extra={**context, "num_pages": len(images)},
        )

        # Create PowerPoint presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        for idx, image in enumerate(images):
            logger.debug(
                f"Adding slide {idx + 1}/{len(images)}",
                extra={**context, "slide_number": idx + 1},
            )

            # Add blank slide
            blank_slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_slide_layout)

            # Save image temporarily
            img_path = os.path.join(tmp_dir, f"page_{idx + 1}.png")
            image.save(img_path, "PNG")

            # Add image to slide
            left = Inches(0)
            top = Inches(0)
            width = prs.slide_width
            height = prs.slide_height

            slide.shapes.add_picture(img_path, left, top, width=width, height=height)

            # Clean up temp image
            os.remove(img_path)

        # Save PowerPoint file
        base_name = Path(safe_filename).stem
        output_filename = f"{base_name}{suffix}.pptx"
        output_path = os.path.join(tmp_dir, output_filename)

        prs.save(output_path)

        output_size = os.path.getsize(output_path)
        logger.info(
            "PDF to PowerPoint conversion completed",
            extra={
                **context,
                "output_path": output_path,
                "output_size": output_size,
                "num_slides": len(images),
            },
        )

        return input_path, output_path

    except Exception as e:
        logger.exception(
            "PDF to PowerPoint conversion failed",
            extra={**context, "error": str(e)},
        )
        raise ConversionError(f"Failed to convert PDF to PowerPoint: {str(e)}") from e
